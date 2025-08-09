import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


# ---------- Minimalist style helpers ----------
TITLE_FONT_NAME = "Calibri Light"
BODY_FONT_NAME = "Calibri"
CODE_FONT_NAME = "Consolas"
TITLE_FONT_SIZE = Pt(40)
BODY_FONT_SIZE = Pt(20)
ACCENT_COLOR = RGBColor(0x22, 0x55, 0x88)


def style_title(shape):
    title_tf = shape.text_frame
    for p in title_tf.paragraphs:
        for r in p.runs:
            r.font.name = TITLE_FONT_NAME
            r.font.size = TITLE_FONT_SIZE
            r.font.bold = True
            r.font.color.rgb = ACCENT_COLOR
    p = title_tf.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def style_body_textframe(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = BODY_FONT_NAME
            r.font.size = BODY_FONT_SIZE
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)


# ---------- Slide builders ----------

def add_bulleted_slide(prs: Presentation, title: str, bullets: list[str], notes: str | None = None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    # Title
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    # Bullets
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for idx, line in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
    style_body_textframe(tf)

    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.clear()
        notes_tf.text = notes

    return slide


def add_title_slide(prs: Presentation, title: str, subtitle: str, notes: str | None = None):
    layout = prs.slide_layouts[0]  # Title
    slide = prs.slides.add_slide(layout)

    # Title
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    # Subtitle
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
        style_body_textframe(slide.placeholders[1].text_frame)

    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.clear()
        notes_tf.text = notes
    return slide


def add_code_box(slide, left_in, top_in, width_in, height_in, code_text: str):
    tx = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.name = CODE_FONT_NAME
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    return tx


def add_image(slide, image_path: str, left_in, top_in, width_in=None, height_in=None):
    if width_in is not None:
        pic = slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in), width=Inches(width_in))
    elif height_in is not None:
        pic = slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in), height=Inches(height_in))
    else:
        pic = slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in))
    return pic


def add_sidebar_bullets(slide, left_in, top_in, width_in, height_in, title: str, items: list[str]):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.clear()
    # Title line
    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = title
    run0.font.name = TITLE_FONT_NAME
    run0.font.size = Pt(20)
    run0.font.bold = True
    run0.font.color.rgb = ACCENT_COLOR
    p0.space_after = Pt(6)
    # Items
    for it in items:
        p = tf.add_paragraph()
        p.text = it
        p.level = 0
        for r in p.runs:
            r.font.name = BODY_FONT_NAME
            r.font.size = Pt(16)
    return box


def add_budget_pie_chart(slide):
    # Data from proposal
    labels = [
        "Research Stipend",
        "Compute",
        "Software & Tools",
        "Dissemination",
        "Contingency",
    ]
    values = [30000, 17500, 5300, 2500, 1000]

    chart_data = ChartData()
    chart_data.categories = labels
    chart_data.add_series("Budget ($56,300)", values)

    x, y, cx, cy = Inches(6.5), Inches(1.8), Inches(6.0), Inches(3.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.RIGHT

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = "$#,##0"
    data_labels.position = XL_LABEL_POSITION.BEST_FIT

    return chart


def add_approach_diagram(slide):
    # Simple blocks and arrows: Sensors -> Encoders -> HVAE -> Counterfactuals -> Decoder
    left = 0.5
    top = 2.0
    w = 2.1
    h = 1.0

    def add_block(x, label):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(h))  # 1=RECTANGLE
        tf = shape.text_frame
        tf.text = label
        style_body_textframe(tf)
        for p in tf.paragraphs:
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        return shape

    def add_arrow(x_from, x_to):
        y = top + h / 2
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x_from + w), Inches(y), Inches(x_to), Inches(y))  # 2=STRAIGHT

    x_positions = [left, left + 2.6, left + 5.2, left + 7.8, left + 10.4]
    labels = ["Sensors", "Encoders", "HVAE", "Counterfactuals", "Decoder"]
    blocks = []
    for x, lab in zip(x_positions, labels):
        blocks.append(add_block(x, lab))
    for i in range(len(x_positions) - 1):
        add_arrow(x_positions[i], x_positions[i + 1])


def add_novelty_table(slide):
    left, top, width, height = Inches(6.5), Inches(1.4), Inches(6.0), Inches(3.5)
    rows, cols = 6, 4
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Feature", "Our Model", "PlaNet", "RWM"]
    features = [
        ("Modalities", "Visual, proprio, reward, actions", "Visual only", "Visual, action"),
        ("Hierarchy", "3-tier HVAE", "None (flat)", "2-layer nets"),
        ("Fusion", "Cross-attention + dropout", "Concat/shared", "Concat + LSTM"),
        ("Counterfactuals", "Latent trajectory editing", "Rollouts only", "Limited imagination"),
        ("Affect", "Reward valence encoding", "Scalar reward", "Scalar reward"),
    ]

    # Header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = TITLE_FONT_NAME
                r.font.size = Pt(16)
                r.font.bold = True
                r.font.color.rgb = ACCENT_COLOR

    # Body rows
    for r_idx, row in enumerate(features, start=1):
        for c_idx, txt in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = txt
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = BODY_FONT_NAME
                    r.font.size = Pt(14)

    return table_shape


def add_multimodal_mini_diagram(slide):
    left = 6.5
    top = 1.4
    w = 1.8
    h = 0.8
    gap = 0.2

    labels = ["Visual", "Proprio", "Action", "Reward"]
    blocks = []
    for i, lab in enumerate(labels):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + i * (h + gap)), Inches(w), Inches(h))
        tf = shape.text_frame
        tf.text = lab
        style_body_textframe(tf)
        for p in tf.paragraphs:
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        blocks.append(shape)

    # Fusion block
    fx = left + w + 0.6
    fusion = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(fx), Inches(top + 1.0), Inches(2.2), Inches(1.6))
    fusion.text_frame.text = "Cross-Attn\nFusion"
    style_body_textframe(fusion.text_frame)
    for p in fusion.text_frame.paragraphs:
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Connectors
    for i in range(len(blocks)):
        y = top + i * (h + gap) + h / 2
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left + w), Inches(y), Inches(fx), Inches(top + 1.0 + 0.8))

    return fusion


def add_temporal_segmentation_bar(slide):
    left = 6.5
    top = 1.6
    total_w = 6.0
    h = 0.5

    segments = [0.25, 0.35, 0.2, 0.2]
    colors = [RGBColor(0x22, 0x55, 0x88), RGBColor(0xAA, 0x77, 0x33), RGBColor(0x3A, 0x9D, 0x23), RGBColor(0x88, 0x88, 0x88)]

    x = left
    for frac, color in zip(segments, colors):
        w = total_w * frac
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(h))
        fill = rect.fill
        fill.solid()
        fill.fore_color.rgb = color
        rect.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        x += w

    # Label
    box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.6), Inches(total_w), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Episode segmentation (phases)"
    for r in p.runs:
        r.font.name = BODY_FONT_NAME
        r.font.size = Pt(14)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    return box


def add_hierarchy_ladder(slide):
    left = 6.5
    top = 1.4
    w = 2.2
    h = 0.8
    gap = 0.3

    labels = [("z3", "Conceptual"), ("z2", "Segments"), ("z1", "Sensorimotor")]
    for i, (layer, desc) in enumerate(labels):
        y = top + i * (h + gap)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(y), Inches(w), Inches(h))
        tf = rect.text_frame
        tf.text = f"{layer}: {desc}"
        style_body_textframe(tf)
        for p in tf.paragraphs:
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        if i < len(labels) - 1:
            slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left + w / 2), Inches(y + h), Inches(left + w / 2), Inches(y + h + gap))

    return True


def add_gantt_timeline(slide):
    left = 0.5
    top = 4.0
    lane_h = 0.4
    lane_gap = 0.25

    phases = [
        ("Foundation", 0.0, 0.18),
        ("Core Dev", 0.18, 0.35),
        ("Temporal+Hierarchy", 0.35, 0.55),
        ("Eval & Validation", 0.55, 0.75),
        ("Integration & Demo", 0.75, 1.0),
    ]

    total_w = 12.0

    for i, (name, start, end) in enumerate(phases):
        y = top + i * (lane_h + lane_gap)
        x = left + total_w * start
        w = total_w * (end - start)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(lane_h))
        fill = bar.fill
        fill.solid()
        fill.fore_color.rgb = ACCENT_COLOR
        bar.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # label
        tx = slide.shapes.add_textbox(Inches(left + 0.1), Inches(y - 0.05), Inches(5.0), Inches(lane_h))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = name
        for r in p.runs:
            r.font.name = BODY_FONT_NAME
            r.font.size = Pt(14)

    return True


# ---------- Figure generation ----------

def ensure_figures() -> dict:
    """Generate placeholder UMAP-like figure and return paths."""
    out_dir = Path("/workspace/figures")
    out_dir.mkdir(parents=True, exist_ok=True)
    umap_png = out_dir / "umap_example.png"

    if not umap_png.exists():
        try:
            import numpy as np
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt

            rng = np.random.default_rng(42)
            # Three clusters
            centers = np.array([[0, 0], [3, 0.5], [1.5, 2.5]])
            points = []
            colors = []
            for i, c in enumerate(centers):
                pts = rng.normal(loc=c, scale=0.35, size=(200, 2))
                points.append(pts)
                colors.append(np.full(200, i))
            X = np.vstack(points)
            C = np.concatenate(colors)

            plt.figure(figsize=(6, 4))
            palette = ["#225588", "#AA7733", "#3A9D23"]
            for i in range(3):
                sel = C == i
                plt.scatter(X[sel, 0], X[sel, 1], s=10, alpha=0.7, c=palette[i], label=f"Cluster {i+1}")
            plt.title("Latent Trajectory Embeddings (UMAP-style)")
            plt.xticks([])
            plt.yticks([])
            plt.legend(frameon=False)
            plt.tight_layout()
            plt.savefig(umap_png, dpi=150)
            plt.close()
        except Exception:
            pass

    return {"umap": str(umap_png)}


# ---------- Build deck ----------

def build_presentation() -> Presentation:
    prs = Presentation()

    # 1) Title
    add_title_slide(
        prs,
        title="Latent Experience Modeling for Counterfactual Reasoning in Agents",
        subtitle="Presenter: Chris Mangum | CIMC Research Proposal",
        notes=(
            "Set context and credibility.\n"
            "Emphasize: bridging reactive behavior and reflective cognition via latent experience models."
        ),
    )

    # 2) Motivation & CIMC alignment
    slide2 = add_bulleted_slide(
        prs,
        title="Motivation & CIMC Alignment",
        bullets=[
            "Problem: Agents are reactive; lack introspective memory and counterfactual imagination",
            "Fit: Architectures for conscious agents; methodology of modeling consciousness",
            "Goal: From reactive to reflective agents via editable latent memories",
        ],
        notes=(
            "Tie to CIMC focus: computational models of consciousness, self, episodic memory.\n"
            "Visual cue: Reactive → Reflective diagram."
        ),
    )
    # Add clickable link line
    tf2 = slide2.shapes.placeholders[1].text_frame
    p_link = tf2.add_paragraph()
    p_link.text = "CIMC Call for Proposals"
    if p_link.runs:
        p_link.runs[0].hyperlink.address = "https://cimc.ai/#/proposals"
        for r in p_link.runs:
            r.font.name = BODY_FONT_NAME
            r.font.size = Pt(18)
            r.font.color.rgb = ACCENT_COLOR

    # Add right-side CIMC focus areas
    add_sidebar_bullets(
        slide2,
        left_in=6.5,
        top_in=1.4,
        width_in=6.0,
        height_in=3.5,
        title="CIMC focus areas (aligned)",
        items=[
            "Architectures for conscious agents",
            "Methodology of modeling consciousness",
            "Emergence & role of self in learning agents",
            "Environments for intelligent behavior",
        ],
    )

    # 3) Objectives & key questions (add right-side key questions panel)
    slide3 = add_bulleted_slide(
        prs,
        title="Objectives & Key Questions",
        bullets=[
            "Objectives: vectorized episodic memory; smooth latent trajectories; counterfactual editing",
            "Outcome: reusable substrate for introspective reasoning",
        ],
        notes=(
            "Anchor to proposal sections: Research Objectives and Key Research Questions."
        ),
    )
    add_sidebar_bullets(
        slide3,
        left_in=6.5,
        top_in=1.4,
        width_in=6.0,
        height_in=3.5,
        title="Key questions",
        items=[
            "What latent structures preserve semantic + temporal coherence?",
            "How to fuse multimodal signals into aligned representations?",
            "Can we perform meaningful counterfactual transformations?",
            "What diagnostics reveal cognitive usefulness?",
        ],
    )

    # 4) Novelty vs prior work (add table)
    slide4 = add_bulleted_slide(
        prs,
        title="Novelty vs Prior Work",
        bullets=[
            "Cross-attention fusion over simple concatenation",
            "Reward–affect encoding with valence shaping",
            "3-tier HVAE for multi-scale memory",
            "Latent trajectory editing (counterfactuals)",
        ],
        notes=(
            "Comparison: PlaNet, recurrent world models.\n"
            "Stress counterfactual interventions + hierarchy."
        ),
    )
    add_novelty_table(slide4)

    # 5) Approach overview with minimalist diagram
    slide5 = add_bulleted_slide(
        prs,
        title="Approach Overview",
        bullets=[
            "Multimodal encoding",
            "Temporal continuity & segmentation",
            "Hierarchical abstraction",
            "Evaluation & diagnostics",
        ],
        notes=(
            "Show system diagram: sensors → encoders → HVAE → counterfactuals → decoder."
        ),
    )
    add_approach_diagram(slide5)

    # 6) Multimodal encoding (add mini fusion diagram)
    slide6 = add_bulleted_slide(
        prs,
        title="Multimodal Encoding",
        bullets=[
            "Cross-attention fusion; modality dropout",
            "Reward–affect encoder with contrastive valence loss",
            "Full-tuple reconstruction to avoid channel neglect",
        ],
        notes=(
            "Add schematic of cross-attention; show valence separation in latent space."
        ),
    )
    add_multimodal_mini_diagram(slide6)

    # 7) Temporal continuity & segmentation (add segmentation bar)
    slide7 = add_bulleted_slide(
        prs,
        title="Temporal Continuity & Segmentation",
        bullets=[
            "Sequence encoders with smoothness regularization",
            "Surprise-based event boundaries; attention pooling",
            "DTW/Soft-DTW latent trajectory alignment",
        ],
        notes=(
            "Include episode trajectory plot with colored segments."
        ),
    )
    add_temporal_segmentation_bar(slide7)

    # 8) Hierarchical abstraction (add ladder)
    slide8 = add_bulleted_slide(
        prs,
        title="Hierarchical Abstraction",
        bullets=[
            "3-tier HVAE: z1 sensorimotor; z2 segments; z3 conceptual",
            "Attention-based abstraction and auxiliary tasks",
            "Drill-down/roll-up consistency checks",
        ],
        notes=(
            "Ladder VAE diagram; annotate levels and semantics."
        ),
    )
    add_hierarchy_ladder(slide8)

    # 9) Evaluation & diagnostics with UMAP-style figure
    slide9 = add_bulleted_slide(
        prs,
        title="Evaluation & Diagnostics",
        bullets=[
            "Reconstruction, clustering (NMI), temporal coherence",
            "Traversal & perturbation tests; retrieval@K",
            "Counterfactual success: plausibility, human ratings, reward Δ",
        ],
        notes=(
            "Show UMAP plot, metric table, traversal example."
        ),
    )

    figs = ensure_figures()
    if os.path.exists(figs["umap"]):
        add_image(slide9, figs["umap"], left_in=6.5, top_in=1.4, width_in=6.0)

    # 10) Implementation plan & timeline (add Gantt)
    slide10 = add_bulleted_slide(
        prs,
        title="Implementation Plan & Timeline (40 Weeks)",
        bullets=[
            "Phases: Foundations; Core; Temporal+Hierarchy; Eval; Demo",
            "Tools: PyTorch, Optuna, W&B; UMAP/t-SNE",
            "Environment: custom 2D gridworld (AgentFarm-compatible)",
        ],
        notes=(
            "Present as 5-bar Gantt; call out key milestones."
        ),
    )
    add_gantt_timeline(slide10)

    # 11) Deliverables, impact, demo + API snippet
    slide11 = add_bulleted_slide(
        prs,
        title="Deliverables, Impact, and Demo",
        bullets=[
            "Deliverables: model, toolkit, dashboard, dataset, meta-embedding, report",
            "Impact: introspective agents; multimodal fusion advances; evaluation standards",
            "API example shown on right",
        ],
        notes=(
            "Include dashboard mock and simple API snippet."
        ),
    )

    api_snippet = (
        "# Embedding API\n"
        "embedding = model.encode(experience)\n"
        "episodes = model.query(predicate=\"reward>threshold and danger\")\n"
        "edited = model.counterfactual(episode, objective=\"avoid obstacle; reach goal\")"
    )
    add_code_box(slide11, left_in=6.5, top_in=1.4, width_in=6.0, height_in=3.0, code_text=api_snippet)

    # 12) Budget, risks, team + native pie chart
    slide12 = add_bulleted_slide(
        prs,
        title="Budget, Risks, and Team",
        bullets=[
            "Budget summary: stipend, compute, tools, dissemination, contingency",
            "Risks & mitigations: fusion alignment; fidelity vs imagination; generalization; interpretability",
            "Team: prior latent trajectory work; narrative phases",
        ],
        notes=(
            "Pie chart for budget; concise risk table; brief bio."
        ),
    )
    add_budget_pie_chart(slide12)

    # 13) Q&A
    add_bulleted_slide(
        prs,
        title="Q&A",
        bullets=[
            "Evaluation rigor & thresholds",
            "Counterfactual criteria & constraints",
            "Generalization & transfer tests",
            "Dashboard demo plan",
        ],
        notes=(
            "Keep 1–2 minutes buffer; be ready with backup slides."
        ),
    )

    # Backups
    add_bulleted_slide(
        prs,
        title="Backup: Metrics & Thresholds",
        bullets=[
            "SSIM/MSE for recon",
            "NMI/Silhouette for clustering",
            "Curvature/ISTD for smoothness",
            "Retrieval Precision@K",
        ],
        notes="Use only if questioned on evaluation details.",
    )

    add_bulleted_slide(
        prs,
        title="Backup: Counterfactual Editing",
        bullets=[
            "Gradient-based latent optimization",
            "Constraints for plausibility",
            "Success rate & reward delta",
        ],
        notes="Keep formulae in speaker notes if needed.",
    )

    return prs


def main():
    ensure_figures()

    prs = build_presentation()
    output_path = "/workspace/CIMC_Proposal_Presentation.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()